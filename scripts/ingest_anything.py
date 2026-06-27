#!/usr/bin/env python3
"""Universal Smart Knowledge Ingestion.

Drop any textual/knowledge data into knowledge-inbox/ and this script
converts everything to clean Markdown with YAML frontmatter, deduplicates,
tag/classifies, redacts secrets, and writes manifests.

Usage:
  python ingest_anything.py
  python ingest_anything.py --input knowledge-inbox --output security-brain/imported
  python ingest_anything.py --dry-run
  python ingest_anything.py --index
"""

import argparse
import csv
import hashlib
import json
import logging
import mimetypes
import os
import re
import subprocess
import sys
import tempfile
import time
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

import yaml

logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s", datefmt="%H:%M:%S")
log = logging.getLogger("ingest-anything")


# ── Default paths ──────────────────────────────────────────────────
DEFAULT_INPUT = "knowledge-inbox"
DEFAULT_OUTPUT = "security-brain/imported"
DEFAULT_MANIFEST_DIR = "security-brain/manifests"


# ── Config loaders ──────────────────────────────────────────────────
def load_yaml(path: str) -> dict:
    p = Path(path)
    if not p.is_file():
        return {}
    try:
        with open(p) as f:
            return yaml.safe_load(f) or {}
    except Exception:
        return {}


def load_ingestion_config() -> dict:
    cfg = load_yaml("configs/ingestion.yaml")
    if not cfg:
        # Try repo-root relative
        cfg = load_yaml(
            str(Path(__file__).resolve().parent.parent / "configs/ingestion.yaml")
        )
    return cfg


def load_taxonomy() -> dict:
    cfg = load_yaml("configs/content_taxonomy.yaml")
    if not cfg:
        cfg = load_yaml(
            str(Path(__file__).resolve().parent.parent / "configs/content_taxonomy.yaml")
        )
    return cfg.get("topics", {})


# ── Topic detection ─────────────────────────────────────────────────
TOPIC_CACHE: Optional[Dict[str, List[str]]] = None


def _get_topics_dict() -> Dict[str, List[str]]:
    global TOPIC_CACHE
    if TOPIC_CACHE is not None:
        return TOPIC_CACHE
    raw = load_taxonomy()
    out: Dict[str, List[str]] = {}
    for topic_name, topic_data in raw.items():
        if isinstance(topic_data, dict) and "keywords" in topic_data:
            out[topic_name] = [str(k).lower() for k in topic_data["keywords"]]
    TOPIC_CACHE = out
    return out


def detect_topics(text: str) -> List[str]:
    if not text or len(text) < 20:
        return []
    text_lower = text.lower()
    topics = _get_topics_dict()
    matches = []
    for topic, keywords in topics.items():
        score = 0
        for kw in keywords:
            if kw in text_lower:
                score += 1
        if score >= 1:
            matches.append((topic, score))
    matches.sort(key=lambda x: -x[1])
    return [m[0] for m in matches[:6]]


def detect_category(filepath: Path, source_type: str, topics: List[str]) -> str:
    name = (filepath.name or "").lower()
    parent_name = (filepath.parent.name or "").lower() if filepath.parent else ""

    if source_type in ("markdown", "text", "html", "pdf", "docx"):
        if any(kw in name for kw in ("blog", "article", "writeup", "post")):
            return "blogs"
        if any(kw in name for kw in ("note", "journal", "todo", "scratch")):
            return "notes"
        return "documents"

    if source_type in ("json", "csv", "spreadsheet"):
        return "datasets"

    if source_type == "code":
        return "code"

    if source_type == "url":
        return "blogs"

    if parent_name in ("blogs", "medium", "portswigger"):
        return "blogs"
    if parent_name in ("notes", "journal", "scratch"):
        return "notes"
    if parent_name in ("datasets", "data", "csv"):
        return "datasets"
    if parent_name in ("code", "scripts", "src"):
        return "code"

    return "unknown"


# ── Secret redaction ─────────────────────────────────────────────────
SECRET_PATTERNS: List[Tuple[str, str]] = [
    (r'(?i)(api[_-]?key|apikey|api[_-]?secret|secret[_-]?key)\s*[:=]\s*["\']?([^\s"\']{8,})["\']?',
     r'\1=***REDACTED***'),
    (r'(?i)(password|passwd|pwd)\s*[:=]\s*["\']?([^\s"\']+)["\']?',
     r'\1=***REDACTED***'),
    (r'(?i)(token|jwt|bearer)\s+["\']?([A-Za-z0-9\-_\.]{20,})["\']?',
     r'\1 ***REDACTED***'),
    (r'(?i)(aws[_-]?(access[_-]?key|secret)|AKIA[A-Z0-9]{16})',
     '***REDACTED-AWS-KEY***'),
    (r'(?i)(ghp|gho|ghu|ghs|ghr)_[A-Za-z0-9_]{20,}',
     '***REDACTED-GH-TOKEN***'),
    (r'(?i)(-----BEGIN\s(?:RSA\s)?PRIVATE\sKEY-----[^-]*-----END\s(?:RSA\s)?PRIVATE\sKEY-----)',
     '***REDACTED-PRIVATE-KEY***'),
    (r'(?i)(sk-[A-Za-z0-9_/]{20,})',
     '***REDACTED-API-KEY***'),
    (r'(?i)([a-z0-9]{32,40})\s*(?:\n|$)', '***REDACTED-SUSPECT-TOKEN***'),
]


def redact_secrets(text: str) -> Tuple[str, bool]:
    changed = False
    result = text
    for pattern, replacement in SECRET_PATTERNS:
        new_result, count = re.subn(pattern, replacement, result, flags=re.MULTILINE)
        if count > 0:
            changed = True
            result = new_result
    return result, changed


# ── Deduplication ────────────────────────────────────────────────────
def load_fingerprints(manifest_dir: Path) -> Dict[str, str]:
    fp_file = manifest_dir / "content_fingerprints.json"
    if fp_file.is_file():
        try:
            return json.loads(fp_file.read_text(encoding="utf-8"))
        except Exception:
            return {}
    return {}


def save_fingerprints(manifest_dir: Path, fps: Dict[str, str]) -> None:
    manifest_dir.mkdir(parents=True, exist_ok=True)
    (manifest_dir / "content_fingerprints.json").write_text(
        json.dumps(fps, indent=2, sort_keys=True), encoding="utf-8"
    )


def is_duplicate(text: str, raw: bytes, fingerprints: Dict[str, str]) -> Tuple[bool, str, str]:
    raw_hash = hashlib.sha256(raw).hexdigest()
    text_hash = hashlib.sha256(text.encode("utf-8", errors="replace") or b"").hexdigest()
    if raw_hash in fingerprints or text_hash in fingerprints:
        return True, raw_hash, text_hash
    return False, raw_hash, text_hash


def strip_frontmatter(text: str) -> str:
    if text.startswith("---\n"):
        parts = text.split("---\n", 2)
        if len(parts) >= 3:
            return parts[2].strip()
    return text.strip()


# ── Title detection ──────────────────────────────────────────────────
def detect_title(text: str, filepath: Path) -> str:
    text_stripped = strip_frontmatter(text)
    if not text_stripped:
        return filepath.stem.replace("-", " ").replace("_", " ").title()

    # Check YAML frontmatter first
    if text.startswith("---\n"):
        parts = text.split("---\n", 2)
        if len(parts) >= 2:
            try:
                fm = yaml.safe_load(parts[1])
                if isinstance(fm, dict) and fm.get("title"):
                    return fm["title"]
            except Exception:
                pass

    # First # heading
    m = re.search(r'^#\s+(.+)$', text_stripped, re.MULTILINE)
    if m:
        return m.group(1).strip()

    # HTML title
    m = re.search(r'<title[^>]*>([^<]+)</title>', text_stripped, re.IGNORECASE)
    if m:
        return m.group(1).strip()

    # First non-empty line (shorter of first 2 lines)
    lines = [l.strip() for l in text_stripped.split("\n") if l.strip() and len(l.strip()) > 10]
    if lines:
        candidate = lines[0]
        if len(candidate) > 120:
            candidate = candidate[:117] + "..."
        return candidate

    return filepath.stem.replace("-", " ").replace("_", " ").title()


# ── Text extraction parsers (graceful degradation) ──────────────────
def extract_pdf(filepath: Path) -> Optional[str]:
    try:
        import fitz
        doc = fitz.open(str(filepath))
        text = "\n".join(page.get_text() for page in doc)
        doc.close()
        return text.strip() if text.strip() else None
    except ImportError:
        pass
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(filepath))
        text = "\n".join(page.extract_text() or "" for page in reader.pages)
        return text.strip() if text.strip() else None
    except ImportError:
        return None
    except Exception:
        return None


def extract_docx(filepath: Path) -> Optional[str]:
    try:
        from docx import Document as DocxDocument
        doc = DocxDocument(str(filepath))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except ImportError:
        return None
    except Exception:
        return None


def extract_pptx(filepath: Path) -> Optional[str]:
    try:
        from pptx import Presentation
        prs = Presentation(str(filepath))
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    lines.extend(
                        rt.text for para in shape.text_frame.paragraphs
                        for rt in para.runs if rt.text
                    )
        return "\n".join(lines) if lines else None
    except ImportError:
        return None
    except Exception:
        return None


def extract_html(filepath: Path) -> Optional[str]:
    raw = filepath.read_text(encoding="utf-8", errors="replace")
    if not raw.strip():
        return None
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(raw, "html.parser")
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        return soup.get_text(separator="\n", strip=True)
    except ImportError:
        # Fallback: strip HTML tags with regex
        stripped = re.sub(r'<script[^>]*>.*?</script>', '', raw, flags=re.DOTALL | re.IGNORECASE)
        stripped = re.sub(r'<style[^>]*>.*?</style>', '', stripped, flags=re.DOTALL | re.IGNORECASE)
        stripped = re.sub(r'<[^>]+>', ' ', stripped)
        stripped = re.sub(r'&[a-z]+;', ' ', stripped)
        stripped = re.sub(r'\s+', ' ', stripped).strip()
        return stripped if len(stripped) > 20 else None
    except Exception:
        return None


def extract_spreadsheet(filepath: Path) -> Optional[str]:
    try:
        import pandas as pd
        ext = filepath.suffix.lower()
        if ext in (".csv", ".tsv"):
            sep = "\t" if ext == ".tsv" else ","
            df = pd.read_csv(str(filepath), sep=sep, nrows=1000)
        elif ext in (".xlsx", ".xls"):
            df = pd.read_excel(str(filepath), nrows=1000)
        else:
            return None
        lines = [",".join(str(c) for c in df.columns)]
        for _, row in df.iterrows():
            lines.append(",".join(str(v) for v in row))
        return "\n".join(lines)
    except ImportError:
        pass
    except Exception:
        pass

    # CSV stdlib fallback
    try:
        with open(filepath, encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            lines = []
            for i, row in enumerate(reader):
                if i > 1000:
                    break
                lines.append(",".join(row))
        return "\n".join(lines) if lines else None
    except Exception:
        return None


def extract_json(filepath: Path) -> Optional[str]:
    try:
        raw = filepath.read_text(encoding="utf-8", errors="replace")
        if filepath.suffix.lower() == ".jsonl":
            lines = []
            for line in raw.strip().split("\n"):
                line = line.strip()
                if not line:
                    continue
                try:
                    obj = json.loads(line)
                    lines.append(json.dumps(obj, indent=2))
                except json.JSONDecodeError:
                    lines.append(line)
            return "\n\n".join(lines) if lines else None
        else:
            data = json.loads(raw)
            return _flatten_json(data)
    except Exception:
        return None


def _flatten_json(data, prefix="") -> str:
    lines = []
    if isinstance(data, dict):
        for k, v in data.items():
            full = f"{prefix}{k}" if prefix else k
            if isinstance(v, (dict, list)):
                lines.append(f"{full}:")
                lines.append(_flatten_json(v, f"  {full}."))
            else:
                lines.append(f"{full}: {v}")
    elif isinstance(data, list):
        for i, item in enumerate(data):
            full = f"{prefix}[{i}]"
            if isinstance(item, (dict, list)):
                lines.append(_flatten_json(item, full))
            else:
                lines.append(f"{full}: {item}")
    return "\n".join(lines)


# ── Archive extraction ──────────────────────────────────────────────
def extract_archive(filepath: Path, temp_dir: Path) -> List[Path]:
    extracted: List[Path] = []
    max_files = 500
    max_total_bytes = 200 * 1024 * 1024  # 200MB total
    total_bytes = 0
    try:
        import shutil
        name_lower = filepath.name.lower()
        if filepath.suffix.lower() == ".zip":
            import zipfile
            with zipfile.ZipFile(str(filepath), "r") as zf:
                for name in zf.namelist():
                    # Safety: skip symlinks, absolute paths, path traversal
                    info = zf.getinfo(name)
                    if info.is_symlink():
                        continue
                    if os.path.isabs(name) or ".." in Path(name).parts:
                        log.warning(f"  Skipping unsafe path in archive: {name}")
                        continue
                    info_size = info.file_size
                    if total_bytes + info_size > max_total_bytes:
                        log.warning(f"  Archive extraction size limit reached ({max_total_bytes//1024//1024}MB)")
                        break
                    if len(extracted) >= max_files:
                        log.warning(f"  Archive file count limit reached ({max_files})")
                        break
                    dest = temp_dir / Path(name).name
                    try:
                        zf.extract(name, temp_dir)
                        extracted.append(temp_dir / name)
                        total_bytes += info_size
                    except Exception as e:
                        log.warning(f"  Could not extract {name}: {e}")
        elif ".tar" in name_lower:
            import tarfile
            with tarfile.open(str(filepath), "r:*") as tf:
                for member in tf.getmembers():
                    if member.issym() or member.isdev():
                        continue
                    if os.path.isabs(member.name) or ".." in Path(member.name).parts:
                        log.warning(f"  Skipping unsafe path: {member.name}")
                        continue
                    if total_bytes + (member.size if member.isfile() else 0) > max_total_bytes:
                        log.warning(f"  Archive extraction size limit reached")
                        break
                    if len(extracted) >= max_files:
                        log.warning(f"  Archive file count limit reached")
                        break
                    try:
                        tf.extract(member, temp_dir, filter="data")
                        extracted.append(temp_dir / member.name)
                        total_bytes += member.size if member.isfile() else 0
                    except Exception as e:
                        log.warning(f"  Could not extract {member.name}: {e}")
    except ImportError:
        return []
    except Exception as e:
        log.warning(f"  Archive extraction error: {e}")
        return []
    return extracted


# ── Clean text ───────────────────────────────────────────────────────
def clean_text(text: str) -> str:
    if not text:
        return ""
    text = text.replace("\x00", "")
    text = re.sub(r'\r\n', '\n', text)
    text = re.sub(r'\r', '\n', text)
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r'[ \t]{3,}', '  ', text)
    text = text.strip()
    return text


# ── Build Markdown output ────────────────────────────────────────────
def build_markdown(
    title: str,
    content: str,
    source_type: str,
    original_path: str,
    original_filename: str,
    category: str,
    topics: List[str],
    raw_sha256: str,
    text_sha256: str,
    redacted: bool,
    language: str = "en",
    source_url: str = "",
) -> str:
    tags = list(dict.fromkeys(["imported", category] + topics))

    fm = {
        "source": "imported",
        "source_type": source_type,
        "original_path": original_path,
        "original_filename": original_filename,
        "title": title,
        "category": category,
        "detected_topics": topics,
        "tags": tags,
        "language": language,
        "raw_sha256": raw_sha256,
        "text_sha256": text_sha256,
        "ingested_at": datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ"),
        "sensitivity": "unknown",
        "redactions_applied": redacted,
    }
    if source_url:
        fm["source_url"] = source_url

    fm_yaml = yaml.dump(fm, default_flow_style=False, allow_unicode=True, sort_keys=False)
    md_content = content.replace("```", "'''").strip()

    return f"""---
{fm_yaml}---

# {title}

## Source Metadata

- Original Path: {original_path}
- Source Type: {source_type}
- Detected Topics: {", ".join(topics) if topics else "none"}
- Ingested At: {fm["ingested_at"]}
- Redactions Applied: {redacted}
- Raw SHA256: `{raw_sha256}`
- Text SHA256: `{text_sha256}`
{f"- Source URL: {source_url}" if source_url else ""}

## Content

{md_content}
"""


# ── Determine source type ────────────────────────────────────────────
def resolve_source_type(suffix: str) -> str:
    suffix = suffix.lower()
    for stype, exts in (load_ingestion_config() or {}).get("supported_extensions", {}).items():
        for ext in exts:
            if ext.lower() == suffix:
                return stype
    # Fallbacks
    type_map = {
        ".md": "markdown", ".markdown": "markdown",
        ".txt": "text", ".log": "text",
        ".html": "html", ".htm": "html",
        ".pdf": "pdf",
        ".docx": "docx", ".doc": "docx",
        ".pptx": "pptx",
        ".csv": "spreadsheet", ".tsv": "spreadsheet",
        ".xlsx": "spreadsheet", ".xls": "spreadsheet",
        ".json": "json", ".jsonl": "json",
        ".zip": "archive", ".tar": "archive",
        ".tar.gz": "archive", ".tgz": "archive",
    }
    if suffix in type_map:
        return type_map[suffix]
    if suffix in (".py", ".js", ".ts", ".go", ".rs", ".java", ".rb", ".php",
                  ".c", ".cpp", ".h", ".cs", ".sh", ".yaml", ".yml", ".toml", ".ini"):
        return "code"
    return "unknown"


# ── Is this a URLs file? ────────────────────────────────────────────
def is_urls_file(filepath: Path) -> bool:
    return filepath.suffix.lower() == ".txt" and filepath.name in ("urls.txt", "urls", "url_list.txt")


# ── Appender helpers ─────────────────────────────────────────────────
def append_manifest(manifest_dir: Path, entry: dict) -> None:
    manifest_dir.mkdir(parents=True, exist_ok=True)
    with open(manifest_dir / "ingest_manifest.jsonl", "a", encoding="utf-8") as f:
        f.write(json.dumps(entry, default=str) + "\n")


def append_error(manifest_dir: Path, filepath: str, error: str, stage: str = "parse") -> None:
    manifest_dir.mkdir(parents=True, exist_ok=True)
    with open(manifest_dir / "ingest_errors.jsonl", "a", encoding="utf-8") as f:
        f.write(json.dumps({
            "filepath": filepath,
            "error": str(error),
            "stage": stage,
            "timestamp": datetime.now(timezone.utc).isoformat(),
        }) + "\n")


# ── Make safe filename ───────────────────────────────────────────────
def safe_out_filename(title: str, raw_hash: str, suffix: str = ".md") -> str:
    try:
        from slugify import slugify
        slug = slugify(title, lowercase=True, max_length=100, word_boundary=True)
    except ImportError:
        slug = re.sub(r'[^\w\s-]', '', title.lower())
        slug = re.sub(r'[-\s]+', '-', slug)[:100]
    if not slug:
        slug = "doc"
    return f"{slug}_{raw_hash[:8]}{suffix}"


# ── Sensitive file denylist ─────────────────────────────────────────
SENSITIVE_PATTERNS = {
    ".env", ".env.local", ".env.production", ".env.development",
    ".npmrc", ".pypirc", "id_rsa", "id_ed25519", "id_ecdsa",
    ".git-credentials", "credentials.json", "service-account.json",
    "docker-compose.override.yml",
}
SENSITIVE_SUFFIXES = {".pem", ".key", ".p12", ".pfx", ".jks", ".keystore"}


def is_sensitive(filepath: Path) -> bool:
    name = filepath.name.lower()
    if name in SENSITIVE_PATTERNS:
        return True
    if filepath.suffix.lower() in SENSITIVE_SUFFIXES:
        return True
    if ".kube" in filepath.parts and "config" in filepath.name:
        return True
    if filepath.name.startswith(".") and filepath.name not in (".gitkeep", ".gitignore"):
        return True
    return False
def process_file(
    filepath: Path,
    output_dir: Path,
    manifest_dir: Path,
    fingerprints: Dict[str, str],
    args: argparse.Namespace,
) -> Tuple[int, int, int]:
    """Process a single file. Returns (ingested, skipped, errored)."""
    max_size = (args.max_file_size_mb or 50) * 1024 * 1024
    try:
        fsize = filepath.stat().st_size
    except OSError:
        return 0, 0, 0
    if fsize > max_size:
        append_error(manifest_dir, str(filepath), f"File too large ({fsize//1024//1024}MB > {args.max_file_size_mb}MB)")
        return 0, 0, 1
    if fsize == 0:
        return 0, 0, 0

    # Skip sensitive files unless allowed
    if not args.allow_sensitive and is_sensitive(filepath):
        append_error(manifest_dir, str(filepath), "Sensitive file skipped (use --allow-sensitive to force)")
        return 0, 0, 1

    # Read raw bytes for fingerprinting
    try:
        raw_bytes = filepath.read_bytes()
    except Exception as e:
        append_error(manifest_dir, str(filepath), str(e))
        return 0, 0, 1

    suffix = "".join(filepath.suffixes).lower() if filepath.suffixes else filepath.suffix.lower()
    if suffix.endswith(".tar.gz"):
        suffix = ".tar.gz"
    source_type = resolve_source_type(suffix)
    if source_type == "unknown":
        return 0, 0, 0

    # Extract text based on source type
    raw_text: Optional[str] = None
    parse_error: Optional[str] = None
    original_path = str(filepath)

    if source_type == "markdown":
        raw_text = filepath.read_text(encoding="utf-8", errors="replace")
    elif source_type == "text":
        raw_text = filepath.read_text(encoding="utf-8", errors="replace")
    elif source_type == "html":
        raw_text = extract_html(filepath)
        if raw_text is None:
            parse_error = "HTML extraction failed (install beautifulsoup4 for better results)"
    elif source_type == "pdf":
        raw_text = extract_pdf(filepath)
        if raw_text is None:
            parse_error = "PDF extraction failed (install pymupdf or pypdf)"
    elif source_type == "docx":
        raw_text = extract_docx(filepath)
        if raw_text is None:
            parse_error = "DOCX extraction failed (install python-docx)"
    elif source_type == "pptx":
        raw_text = extract_pptx(filepath)
        if raw_text is None:
            parse_error = "PPTX extraction failed (install python-pptx)"
    elif source_type == "spreadsheet":
        raw_text = extract_spreadsheet(filepath)
        if raw_text is None:
            parse_error = "Spreadsheet extraction failed (install pandas or use CSV)"
    elif source_type == "json":
        raw_text = extract_json(filepath)
        if raw_text is None:
            parse_error = "JSON extraction failed"
    elif source_type == "code":
        raw_code = filepath.read_text(encoding="utf-8", errors="replace")
        lang = suffix.lstrip(".")
        raw_text = f"```{lang}\n{raw_code}\n```"
    elif source_type == "archive":
        if not args.extract_archives:
            return 0, 0, 0
        with tempfile.TemporaryDirectory() as tmpdir:
            tmp = Path(tmpdir)
            extracted = extract_archive(filepath, tmp)
            ingested = 0
            for ex in extracted:
                i, _, _ = process_file(ex, output_dir, manifest_dir, fingerprints, args)
                ingested += i
            return ingested, 0, 0
    else:
        return 0, 0, 0

    if raw_text is None or not raw_text.strip():
        if parse_error:
            append_error(manifest_dir, original_path, parse_error)
        return 0, 0, 1

    # Clean
    cleaned = clean_text(raw_text)
    if not cleaned:
        return 0, 0, 0

    # Deduplicate
    is_dupe, raw_hash, text_hash = is_duplicate(cleaned, raw_bytes, fingerprints)
    if is_dupe and args.dedupe and not args.force:
        return 0, 1, 0

    # Redact secrets
    redacted = False
    if args.redact_secrets:
        cleaned, redacted = redact_secrets(cleaned)

    # Detect metadata
    title = detect_title(cleaned, filepath) if raw_text else filepath.stem
    topics = detect_topics(cleaned)
    category = detect_category(filepath, source_type, topics)

    # Target output directory
    cat_dir = output_dir / category
    cat_dir.mkdir(parents=True, exist_ok=True)

    out_fn = safe_out_filename(title, raw_hash)
    out_path = cat_dir / out_fn

    if out_path.exists() and not args.force:
        return 0, 1, 0

    source_url = ""
    if cleaned.startswith("http"):
        for line in cleaned.split("\n")[:5]:
            if line.strip().startswith("http"):
                source_url = line.strip()
                break

    md = build_markdown(
        title=title,
        content=cleaned,
        source_type=source_type,
        original_path=original_path,
        original_filename=filepath.name,
        category=category,
        topics=topics,
        raw_sha256=raw_hash,
        text_sha256=text_hash,
        redacted=redacted,
        source_url=source_url,
    )

    out_path.write_text(md, encoding="utf-8")

    # Update fingerprints
    fingerprints[raw_hash] = original_path
    fingerprints[text_hash] = original_path

    # Append manifest
    append_manifest(manifest_dir, {
        "source": str(filepath),
        "output": str(out_path),
        "source_type": source_type,
        "category": category,
        "title": title,
        "topics": topics,
        "redacted": redacted,
        "raw_size": fsize,
        "raw_sha256": raw_hash,
        "text_sha256": text_hash,
        "ingested_at": datetime.now(timezone.utc).isoformat(),
    })

    return 1, 0, 0


def run_pipeline(args: argparse.Namespace) -> dict:
    input_dir = Path(args.input)
    output_dir = Path(args.output)
    manifest_dir = Path(args.manifest_dir)

    output_dir.mkdir(parents=True, exist_ok=True)

    if not input_dir.is_dir():
        log.error(f"Input directory does not exist: {input_dir}")
        sys.exit(1)

    fingerprints = load_fingerprints(manifest_dir) if args.dedupe else {}

    stats: Dict[str, Any] = {
        "scanned": 0,
        "ingested": 0,
        "skipped_dupes": 0,
        "errors": 0,
        "by_category": {},
        "by_type": {},
    }

    for filepath in sorted(input_dir.rglob("*")):
        if filepath.is_dir():
            continue
        if not args.include_hidden and filepath.name.startswith("."):
            continue
        if not args.include_hidden and any(p.startswith(".") for p in filepath.parts):
            continue
        if filepath.name == ".gitkeep":
            continue

        stats["scanned"] += 1

        # Handle urls.txt separately
        if is_urls_file(filepath):
            log.info(f"URLs file detected: {filepath.name} — use scripts/ingest_urls.py directly")
            continue

        if args.dry_run:
            suffix = filepath.suffix.lower()
            stype = resolve_source_type(suffix)
            is_sens = is_sensitive(filepath) and not args.allow_sensitive
            if is_sens:
                log.info(f"  [DRY-RUN] Would SKIP (sensitive): {filepath.name}")
                stats["would_skip_sensitive"] = stats.get("would_skip_sensitive", 0) + 1
            elif stype != "unknown":
                log.info(f"  [DRY-RUN] Would process: {filepath.name} ({stype})")
                stats["would_ingest"] = stats.get("would_ingest", 0) + 1
                stats["by_type_dry"] = stats.get("by_type_dry", {})
                stats["by_type_dry"][stype] = stats["by_type_dry"].get(stype, 0) + 1
            else:
                stats["would_skip_unknown"] = stats.get("would_skip_unknown", 0) + 1
            continue

        ingested, skipped, errored = process_file(
            filepath, output_dir, manifest_dir, fingerprints, args
        )
        stats["ingested"] += ingested
        stats["skipped_dupes"] += skipped
        stats["errors"] += errored

        if ingested:
            suffix = filepath.suffix.lower()
            stype = resolve_source_type(suffix)
            stats["by_category"] = stats.get("by_category", {})
            stats["by_type"][stype] = stats["by_type"].get(stype, 0) + 1
            cat = detect_category(filepath, stype, [])
            stats["by_category"][cat] = stats["by_category"].get(cat, 0) + 1

    # Save fingerprints
    if args.dedupe and not args.dry_run:
        save_fingerprints(manifest_dir, fingerprints)

    return stats


def main():
    parser = argparse.ArgumentParser(description="Universal Smart Knowledge Ingestion")
    parser.add_argument("--input", default=DEFAULT_INPUT, help=f"Input directory (default: {DEFAULT_INPUT})")
    parser.add_argument("--output", default=DEFAULT_OUTPUT, help=f"Output directory (default: {DEFAULT_OUTPUT})")
    parser.add_argument("--manifest-dir", default=DEFAULT_MANIFEST_DIR, help=f"Manifest directory (default: {DEFAULT_MANIFEST_DIR})")
    parser.add_argument("--profile", default="security", help="Ingestion profile")
    parser.add_argument("--dedupe", action="store_true", default=True, help="Enable deduplication (default: True)")
    parser.add_argument("--no-dedupe", dest="dedupe", action="store_false", help="Disable deduplication")
    parser.add_argument("--force", action="store_true", help="Force re-processing even if already ingested")
    parser.add_argument("--dry-run", action="store_true", help="Scan and report without writing files")
    parser.add_argument("--redact-secrets", action="store_true", default=True, help="Redact secrets (default: True)")
    parser.add_argument("--no-redact", dest="redact_secrets", action="store_false", help="Disable secret redaction")
    parser.add_argument("--extract-archives", action="store_true", help="Extract and process archive files")
    parser.add_argument("--index", action="store_true", help="Trigger knowledge-rag reindex after ingestion")
    parser.add_argument("--max-file-size-mb", type=int, default=50, help="Max file size in MB (default: 50)")
    parser.add_argument("--include-hidden", action="store_true", help="Include hidden files and directories")
    parser.add_argument("--allow-sensitive", action="store_true", help="Allow ingestion of sensitive files (.env, .key, .pem, etc.)")
    parser.add_argument("--verbose", action="store_true", help="Verbose logging")

    args = parser.parse_args()

    if args.verbose:
        log.setLevel(logging.DEBUG)

    print()
    print("=" * 60)
    print("  Universal Smart Knowledge Ingestion")
    print("=" * 60)
    print(f"  Input:       {args.input}")
    print(f"  Output:      {args.output}")
    print(f"  Manifests:   {args.manifest_dir}")
    print(f"  Dedupe:      {args.dedupe}")
    print(f"  Redact:      {args.redact_secrets}")
    print(f"  Extract:     {args.extract_archives}")
    print(f"  Index:       {args.index}")
    print(f"  Dry-run:     {args.dry_run}")
    print()

    if args.dry_run:
        log.info("DRY RUN — no files will be written")

    stats = run_pipeline(args)
    sc = stats["scanned"]
    ig = stats["ingested"]
    du = stats["skipped_dupes"]
    er = stats["errors"]

    print()
    print("-" * 40)
    print(f"  Files scanned:        {sc}")
    if args.dry_run:
        print(f"  Would ingest:         {stats.get('would_ingest', 0)}")
        print(f"  Would skip (sensitive): {stats.get('would_skip_sensitive', 0)}")
        print(f"  Would skip (unknown):  {stats.get('would_skip_unknown', 0)}")
        if stats.get("by_type_dry"):
            print(f"  By type:              {json.dumps(stats['by_type_dry'])}")
    else:
        print(f"  Files ingested:       {ig}")
        print(f"  Duplicates skipped:   {du}")
        print(f"  Errors:               {er}")
        if stats.get("by_type"):
            print(f"  By type:              {json.dumps(stats['by_type'])}")
        if stats.get("by_category"):
            print(f"  By category:          {json.dumps(stats['by_category'])}")
    print("-" * 40)

    if not args.dry_run:
        log.info(f"Output written to {args.output}")
        manifest_file = Path(args.manifest_dir) / "ingest_manifest.jsonl"
        errors_file = Path(args.manifest_dir) / "ingest_errors.jsonl"
        if manifest_file.exists():
            log.info(f"Manifest: {manifest_file}")
        if errors_file.exists():
            log.info(f"Errors:   {errors_file}")

    # Optional indexing
    if args.index and not args.dry_run and ig > 0:
        log.info("Triggering reindex...")
        index_script = Path(__file__).resolve().parent / "reindex_knowledge.py"
        if index_script.is_file():
            subprocess.run(
                [sys.executable, str(index_script), "--path", args.output],
                check=False,
            )
        else:
            log.warning("reindex_knowledge.py not found — reindex manually")

    success = er == 0
    return 0 if success else 1


if __name__ == "__main__":
    sys.exit(main())
